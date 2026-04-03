import streamlit as st
import os
import re
import requests
import json
import time
import pandas as pd
import PyPDF2
import docx
import pptx
import urllib3
from io import BytesIO
from datetime import datetime
from fpdf import FPDF
from langchain_anthropic import ChatAnthropic
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.messages import SystemMessage, HumanMessage

# ==========================================
# 🚨 [추가된 코드] 사내 네트워크 SSL 인증서 에러 해결 🚨
# 이 코드는 다른 모듈들이 호출되기 전에 가장 먼저 실행되어야 합니다.
import ssl
import httpx

# 1. 파이썬 기본 SSL 검증 무력화 (권장하지 않지만 로컬 테스트용으로 가장 빠름)
ssl._create_default_https_context = ssl._create_unverified_context
# SSL 인증서 우회 시 발생하는 터미널 경고 메시지 숨기기
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

# 2. httpx 통신 시 인증서 검증 무력화 환경변수 설정
os.environ['CURL_CA_BUNDLE'] = ''
os.environ['REQUESTS_CA_BUNDLE'] = ''
os.environ['SSL_CERT_FILE'] = ''
os.environ["HTTPX_SSL_VERIFY"] = "0"
# ==========================================

# --- 1. API 키 설정 ---
ANT_KEY = st.secrets.get("ANTHROPIC_API_KEY")
GEM_KEY = st.secrets.get("GOOGLE_API_KEY")
MANUS_API_KEY = st.secrets.get("MANUS_API_KEY")

# LLM 초기화 (모델명 고정)
fast_llm = ChatGoogleGenerativeAI(model="gemini-3-flash-preview", google_api_key=GEM_KEY)
smart_llm = ChatAnthropic(model="claude-sonnet-4-6", anthropic_api_key=ANT_KEY)

# --- 2. 초기 세션 상태 관리 ---
if "step" not in st.session_state:
    st.session_state.step = 1
for key in ["draft_content", "teaser_content", "design_recommendation", "month_title"]:
    if key not in st.session_state:
        st.session_state[key] = ""
if "manus_status" not in st.session_state:
    st.session_state.manus_status = "idle"
if "manus_url" not in st.session_state:
    st.session_state.manus_url = None

# --- 3. 유틸리티 함수 ---
def extract_text_from_file(uploaded_file):
    text = ""
    file_extension = uploaded_file.name.split(".")[-1].lower()
    try:
        if file_extension == "pdf":
            pdf_reader = PyPDF2.PdfReader(uploaded_file)
            for page in pdf_reader.pages:
                extracted = page.extract_text()
                if extracted: text += extracted + "\n"
        elif file_extension in ["docx", "doc"]:
            doc = docx.Document(uploaded_file)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif file_extension in ["pptx", "ppt"]:
            ppt = pptx.Presentation(uploaded_file)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif file_extension in ["xlsx", "xls"]:
            df = pd.read_excel(uploaded_file)
            text += df.to_string() + "\n"
    except Exception as e:
        text += f"\n[{uploaded_file.name} 파일 읽기 오류: {e}]\n"
    return text

def create_manus_infographic(topic, report_content, style_instruction):
    if not MANUS_API_KEY: return None, "Manus API 키 누락"
    url = "https://api.manus.ai/v1/tasks"
    headers = {"API_KEY": MANUS_API_KEY, "Content-Type": "application/json"}
    
    enhanced_prompt = f"""
    반드시 적은 스타일대로 슬라이드를 생성하고 이 스타일을 벗어나서는 안돼. 
    페이지별로 어떤 정보가 들어가야 하는지도 적어 놨으니 이를 참고해서 슬라이드를 구상하고 생성해줘. 
    슬라이드 내 텍스트는 너무 길지 않아야 하고 핵심 내용 중심으로 작성해줘. 특히 간지 페이지에는 절대로 임의로 내용을 추가해서는 안돼.
    시각적으로 돋보이는 16:9 비율의 전문적인 인포그래픽 프레젠테이션을 만들어주세요. 

    [디자인 가이드라인]
    ※ 가이드 라인은 절대로 수정하지 말고 그대로 반영
    1. 시각화 중심: 텍스트를 그대로 나열하기보다는, 내용을 한눈에 이해할 수 있도록 다이어그램, 차트, 인포그래픽 등을 적절히 활용해 주세요.
    2. 간결한 텍스트: 긴 문단보다는 핵심 키워드와 짧은 글머리 기호(Bullet points)를 사용하여 내용을 깔끔하게 요약해 주시면 좋습니다.
    3. 풍부한 시각 자료: 슬라이드 내용과 어울리는 고품질 이미지, 관련 아이콘 또는 일러스트를 자연스럽게 배치해 시각적인 매력을 더해 주세요.
    4. 세련된 레이아웃: 여백을 충분히 두고, 가독성이 좋은 폰트를 사용하여 매거진이나 테크 행사 발표 자료처럼 세련되고 깔끔한 느낌을 연출해 주세요.
    5. 아래 디자인 스타일대로만 슬라이드를 작성하고 절대 디자인 스타일을 벗어나지 말아주세요.

    [디자인 스타일]
    {style_instruction}

    [🚨 필수 텍스트 표기 규칙]
    - 브랜드명 고정: 소식지 이름인 "expl'AI'n telink"는 반드시 소문자 바탕에 'AI'만 대문자로 표기해야 합니다.

    본문 내용:
    {report_content[:4000]}
    """
    
    data = {
        "prompt": enhanced_prompt,
        "agentProfile": "manus-1.6-lite",
        "taskMode": "agent",
        "createShareableLink": True
    }
    
    try:
        response = requests.post(url, json=data, headers=headers, verify=False)
        if response.status_code not in [200, 201]: 
            return None, f"API 연결 실패 (상태 코드: {response.status_code})"
            
        res_json = response.json()
        task_id = res_json.get("task_id")
        task_url = res_json.get("task_url") 
        
        if not task_id: return None, "Task ID를 발급받지 못했습니다."

        st.write("⏳ Manus 서버에서 디자인을 생성하고 있습니다. (최대 10분 소요)")
        
        for _ in range(60):
            time.sleep(10)
            res = requests.get(f"{url}/{task_id}", headers={"API_KEY": MANUS_API_KEY}, verify=False).json()
            
            if res.get("status") == "completed":
                files = res.get("files", [])
                pptx_url = next((f.get('url') for f in files if isinstance(f, dict) and f.get('filename', '').endswith('.pptx')), None)
                share_url = res.get("share_url")
                
                final_url = pptx_url or share_url or task_url
                
                if final_url: 
                    return final_url, "성공"
                else: 
                    return "https://manus.ai", "완료 (Manus 홈페이지에서 확인하세요)"
            
            elif res.get("status") in ["failed", "error"]:
                return None, f"Manus API 내부 오류: {res.get('error', '알 수 없는 오류')}"
                
        return task_url or "https://manus.ai", "제작 시간이 길어지고 있습니다. Manus에서 확인하세요."
    except Exception as e: 
        return None, f"시스템 오류: {str(e)}"

def create_professional_pdf(text, title):
    pdf = FPDF()
    pdf.add_page()
    
    import os
    current_dir = os.path.dirname(os.path.abspath(__file__))
    
    eb_font = os.path.join(current_dir, "NanumSquareEB.ttf")
    r_font = os.path.join(current_dir, "NanumSquareR.ttf")
    
    # 🚨 [수정 1] 폰트 에러 방지를 위해 fname= 파라미터를 정확히 명시
    if os.path.exists(eb_font):
        pdf.add_font('NS_EB', style='', fname=eb_font)
        t_f = 'NS_EB'
    else:
        t_f = 'Helvetica'
        
    if os.path.exists(r_font):
        pdf.add_font('NS_R', style='', fname=r_font)
        b_f = 'NS_R'
    else:
        b_f = 'Helvetica'
    
    clean_title = title
    clean_text = text.replace('#', '').replace('*', '')
    
    pdf.set_font(t_f, size=20)
    pdf.cell(0, 15, txt=clean_title, ln=1, align='L')
    pdf.line(10, pdf.get_y(), 200, pdf.get_y())
    pdf.ln(10)
    
    pdf.set_font(b_f, size=11)
    pdf.multi_cell(0, 8, txt=clean_text)
    
    # 🚨 [수정 2] 스트림릿 다운로드 에러(bytearray)를 막기 위해 bytes()로 감싸기
    pdf_bytes = bytes(pdf.output())

    import re
    return pdf_bytes, re.sub(r'[\\/*?:"<>|]', "", clean_title)

# --- 4. 에이전트 프롬프트 로직 ---
def extract_clean_text(content):
    if isinstance(content, list):
        return "".join([item.get("text", "") for item in content if isinstance(item, dict) and "text" in item])
    return str(content)

def generate_draft(data):
    # AI 보완 여부 텍스트
    def get_ai_inst(flag):
        return "AI 보완 적용하여 내용 확장" if flag else "텍스트 내용의 의미와 수치를 절대 변경하지 말고 그대로 배치할 것"
    
    # 슬라이드 지정 장수 텍스트
    def get_cnt_inst(cnt):
        return f"반드시 {cnt}장의 슬라이드로 분할하여 구성할 것" if cnt else "내용량에 맞춰 AI가 자율적으로 적절한 수의 슬라이드로 분할할 것"

    inst_fin = f"[보완] {get_ai_inst(data['ai_fin'])} / [분량] {get_cnt_inst(data['cnt_fin'])}"
    inst_int = f"[보완] {get_ai_inst(data['ai_int'])} / [분량] {get_cnt_inst(data['cnt_int'])}"
    inst_insight = f"[보완] {get_ai_inst(data['ai_ins'])} / [분량] {get_cnt_inst(data['cnt_ins'])}"
    inst_news = f"[보완] {get_ai_inst(data['ai_news_chk'])} / [분량] {get_cnt_inst(data['cnt_news'])}"

    prompt = f"""
    당신은 사내 소식지 "expl'AI'n telink: 텔링크를 말하다."의 전문 편집장입니다.
    이 출력물은 최종 보고서가 아니라, PPT 자동 생성기에 자동으로 전달될 '슬라이드 스크립트(Outline)'입니다.
    
    [🚨 필수 규칙 및 금지 사항]
    1. 브랜드명 고정: "expl'AI'n telink"는 반드시 소문자 바탕에 'AI'만 대문자로 표기하세요.
    2. 명칭 통일: 모든 단위는 '페이지'가 아닌 '슬라이드'로 통일하세요.
    3. 순차적 번호 부여 (절대 규칙): 슬라이드 번호는 반드시 1부터 시작하는 순차적인 '숫자'로만 작성하세요. (알파벳 A, B, X, Y 등은 절대 사용 금지). 사용자가 지정한 분량에 따라 슬라이드 장수가 늘어나면 3, 4, 5... 순서대로 번호를 계속 올리세요.
    4. 부가 설명 금지: 사용자에게 안내하는 문구나 인사말은 절대 출력하지 마세요.
    5. 형식: 화려한 마크다운 서식이나 표를 만들지 말고, 핵심 텍스트만 간결하게 작성하세요.

    [입력 정보 및 개별 지침]
    - 발행월: {data['month']}
    - 사내 실적: {data['financial']} 
      (지침: {inst_fin})
    - 사내 주요 소식: {data['internal']} 
      (지침: {inst_int})
    - AI Insight (방향성 및 참고 자료): {data['ai_insight']} 
      (지침: {inst_insight})
    - AI 뉴스: {data['ai_news']} 
      (지침: {inst_news})

    [출력 양식]
    ※ 아래 항목들을 순서대로 배치하되, 번호는 반드시 끊기지 않는 순차적 숫자(1, 2, 3, 4, 5...)로 매기세요.

    슬라이드 1. 표지
    - 제목: expl'AI'n telink: 텔링크를 말하다. - {data['month']}
    - 안내 문구: ※ 이 문서는 AI 및 자동화 솔루션으로 제작 되었습니다.
    - 목차: 사내 소식 / AI Insight / AI 뉴스

    슬라이드 2. 간지 (사내 소식)
    - 텍스트: 사내 소식 (재무성과 및 사내 주요 소식)

    슬라이드 3. 재무성과 (※ 지침에 따라 분할 시 슬라이드 4, 5...로 계속 이어서 번호 부여)
    - (실적 텍스트 작성)
    ※ 숫자를 강조하고 크기를 크게 표기, 연관된 이미지 추가

    슬라이드 [이전 번호에 이어서]. 사내 주요 소식
    - (텍스트 작성)
    ※ 본부별 내용과 연관된 이미지 또는 아이콘 추가
    
    슬라이드 [이전 번호에 이어서]. 간지 (AI Insight)
    - 텍스트: AI Insight

    슬라이드 [이전 번호에 이어서]. AI Insight 본문
    - (텍스트 작성)

    슬라이드 [이전 번호에 이어서]. 간지 (AI 뉴스)
    - 텍스트: AI 뉴스

    슬라이드 [이전 번호에 이어서]. AI 뉴스
    - (텍스트 작성)
    """
    
    res = smart_llm.invoke(prompt)
    return extract_clean_text(res.content)

def generate_teaser(data):
    month = data.get("month", "이번 달")
    internal_summary = data.get("internal", "이번 달 사내 소식 없음")
    ai_insight_summary = data.get("ai_insight", "이번 달 AI 인사이트 없음")
    ai_news_summary = data.get("ai_news", "이번 달 AI 뉴스 없음")
    financial_summary = data.get("financial", "실적 데이터 없음")

    # 🚨 [수정] 티저 기획 원칙에 '구도 중심', 'AI 비중 80% 이상', '테마 제안 금지'를 명시했습니다.
    system_prompt = """
    당신은 'expl'AI'n telink' 사내 소식지의 크리에이티브 디렉터입니다.
    당신의 임무는 구성원들의 호기심을 자극할 티저 이미지의 [전체적인 화면 구도와 오브젝트 배치]만을 기획하는 것입니다.

    [티저 기획 4대 원칙 - 엄수]
    1. 스타일 제안 금지: 뒤에서 디자인 테마가 별도로 결정되므로 "수채화풍, 3D, 사이버펑크" 같은 아트 스타일이나 테마는 절대 지정하지 마세요. 오직 '어떤 피사체가 어디에 있는지' 구도만 묘사합니다.
    2. AI 압도적 비중 (80% 이상): 이미지의 메인 피사체와 전체적인 흐름은 무조건 'AI Insight'와 'AI 뉴스'가 주도해야 합니다. 화면 중앙을 꽉 채우는 거대한 메인 메타포로 설정하세요.
    3. 배경 처리 (20% 미만): '사내 실적'과 '사내 소식'은 메인 AI 피사체 주변을 맴도는 작은 아이콘, 희미한 배경 그래픽, 떠다니는 단어 조각 등 아주 작은 '거드는 수준'으로만 배치하세요.
    4. 캠페인 제안 금지: "앞으로 무엇을 하자"는 식의 문구나 묘사는 금지합니다.
    """

    user_prompt = f"""
    [이번 달 소식지 핵심 데이터: {month}]

    1. 🤖 AI Insight (메인 테마 - 비중 80%):
    {ai_insight_summary[:400]}...

    2. 🌍 AI 뉴스 (메인 서포트):
    {ai_news_summary[:300]}...

    3. 🏢 사내 주요 소식 및 실적 (배경 요소 - 비중 20%):
    - 실적: {financial_summary[:200]}...
    - 소식: {internal_summary[:200]}...

    위 내용을 바탕으로 {month} 호기심 자극용 티저 화면의 [구도와 오브젝트 배치안]을 기획해 주세요. 메타포(은유)를 활용하되, 비율 불균형(AI 중심)을 반드시 지켜주세요.
    """

    try:
        messages = [
            SystemMessage(content=system_prompt),
            HumanMessage(content=user_prompt)
        ]
        res = smart_llm.invoke(messages)
        teaser_concept = res.content
    except Exception as e:
        teaser_concept = f"🚨 티저 컨셉 생성 중 오류가 발생했습니다: {e}"

    return teaser_concept

def regenerate_teaser_from_draft(draft_text):
    # 🚨 [수정] 재생성 시에도 AI 중심의 비중과 스타일 제안 금지 규칙을 동일하게 적용합니다.
    prompt = f"""
    사내 소식지 "expl'AI'n telink"의 작성된 전체 스크립트를 바탕으로 티저 기획안을 새로 구상해 주세요. 
    
    [🚨 필수 조건 및 금지 사항]
    1. 스타일 제안 금지: 특정 디자인 테마나 화풍을 지정하지 마세요. 오직 화면의 '전체적인 구도'와 '오브젝트 배치'만 기획합니다.
    2. AI 압도적 비중: 스크립트 중 'AI Insight'와 'AI 뉴스' 내용이 화면의 80% 이상을 차지하는 거대한 메인 피사체가 되어야 합니다.
    3. 보조 요소 배경화: 사내 실적이나 다른 소식들은 메인 피사체 주변의 작은 아이콘, 그래픽 조각 등으로 아주 미미하게 거들도록만 배치하세요.
    4. 표(Table) 작성 금지 / 인사말 금지

    [참고 스크립트(Outline)]
    {draft_text[:4000]}

    [출력 양식 예시]
    ■ 티저 비주얼 구도 기획 (※ 테마 지정 금지)
    - 메인 피사체 (AI Insight 중심): [설명]
    - 보조/배경 요소 (사내소식/실적 중심): [설명]

    ■ 핵심 후킹 카피 (Headline)
    - "[구성원의 시선을 사로잡는 강렬한 한 줄 카피]"
    """
    
    try:
        res = fast_llm.invoke(prompt)
        return extract_clean_text(res.content)
    except Exception as e:
        return f"🚨 티저 재구상 중 오류 발생: {e}"

def revise_draft(current_draft, feedback):
    prompt = f"다음 초안을 사용자 피드백에 맞게 수정하되, PPT 스크립트(Outline) 형태의 간결한 구조는 반드시 유지해.\n\n[초안]\n{current_draft}\n\n[피드백]\n{feedback}"
    res = smart_llm.invoke(prompt)
    return extract_clean_text(res.content)

# 💡 [핵심] 사용자가 원하는 고퀄리티 프롬프트 구조 강제 적용
# 💡 [핵심] 촌스러운 디자인 배제 및 고퀄리티 프롬프트 강제 적용
def get_design_recommendation(month, ai_insight):
    prompt = f"""
    당신은 세계 최고의 프레젠테이션 아트 디렉터입니다.
    발행월({month})과 주요 다룰 내용(AI, 기술)을 바탕으로, **절대 뻔하지 않고 트렌디한** PPT 디자인 테마 3가지를 영문 프롬프트로 기획하세요.

    [🚨 엄격한 금지 사항 (Anti-Patterns)]
    - "Corporate", "Standard business", "Boring", "Basic blue tech", "Clean minimalist" 등 흔하고 지루한 비즈니스 템플릿 스타일은 절대 금지합니다.
    - 추상적이고 모호한 표현(예: "beautiful design", "modern look")을 쓰지 마세요.

    [✅ 필수 디자인 수준 및 포맷]
    - 사용자는 평범한 PPT가 아니라 매거진, 포스터, 감각적인 시각 자료를 원합니다. (예: Brutalism, Neo-Memphis, Cyberpunk, Editorial Zine, Claymorphism, Hand-drawn doodle 등)
    - 반드시 아래 <우수 예시>와 **동일한 구조(Design style, Layout, Color palette)와 디테일한 묘사**를 꽉꽉 채워야 합니다.

    <우수 예시>
    A hand-drawn doodle style infographic slide.

    Design style:
    - black marker sketch illustrations
    - yellow highlighter highlights
    - white crumpled paper background
    - casual whiteboard doodle drawing
    - simple stick figures and sketch icons
    - arrows, stars, speech bubbles, hand drawn shapes
    - handwritten typography style
    - playful startup presentation design

    Layout:
    - vertical infographic
    - two-column split comparison layout
    - section boxes with hand drawn borders
    - circled highlights and underlines
    - sketch diagrams and icons explaining ideas

    Color palette:
    black ink drawing + yellow highlight only
    minimalist visual storytelling infographic
    </우수 예시>

    주요 내용: {ai_insight[:300]}

    출력 형식:
    ### 1. [테마명 (예: Editorial Zine Style)]
    - **추천 이유**: [이유]
    - **프롬프트**:
    ```
    [여기에 예시 구조와 같이 작성된 상세 영문 프롬프트]
    ```
    - **참고**: [🎨 디자인 느낌 미리보기](https://www.google.com/search?tbm=isch&q=Editorial+Zine+presentation+design)

    ### 2. [테마명]
    (동일 구조)
    
    ### 3. [테마명]
    (동일 구조)
    """
    res = fast_llm.invoke(prompt)
    return extract_clean_text(res.content)

# --- 5. UI 영역 ---
st.set_page_config(page_title="expl'AI'n telink Studio", layout="wide")
st.title("📰 expl'AI'n telink 자동화 Agent")

# ---------------------------------------------------------
# Step 1: 정보 입력
# ---------------------------------------------------------
if st.session_state.step == 1:
    st.subheader("📝 Step 1. 이번 달 expl'AI'n telink 데이터 입력")
    
    # 파일 업로더 CSS (유지)
    st.markdown("""
        <style>
        [data-testid="stFileUploader"] section {
            padding: 10px;
        }
        [data-testid="stFileUploader"] {
            min-height: 50px;
        }
        </style>
    """, unsafe_allow_html=True)

    now = datetime.now()
    default_month_str = f"{now.year}년 {now.month}월호"
    
    with st.form("input_form"):
        month = st.text_input("발행 월 (예: 2026년 3월호)", value=default_month_str)
        st.divider()
        
        # 🚨 [수정] 'AI 자율'을 리스트의 맨 앞으로 옮김
        slide_options = ["AI 자율"] + [f"{i}장" for i in range(1, 11)]
        
        # --- 상단 (Row 1): 텔링크 사내 소식 ---
        st.markdown("#### 🏢 텔링크 사내 소식")
        row1_col1, row1_col2 = st.columns(2)
        
        with row1_col1:
            financial = st.text_area("1. 사내 실적 (매출/영업이익 등)", height=150)
            c1, c2 = st.columns([1, 1])
            with c1: 
                st.markdown("<div style='margin-top: 30px;'></div>", unsafe_allow_html=True)
                ai_fin = st.checkbox("✨ AI 보완", key="chk_fin", value=False)
            with c2: 
                # 🚨 순서가 바뀌었으므로 "1장"은 index 1이 됩니다.
                sel_fin = st.selectbox("슬라이드 배분", slide_options, index=1, key="sel_fin")
                
        with row1_col2:
            internal = st.text_area("2. 사내 주요 소식", height=150)
            c3, c4 = st.columns([1, 1])
            with c3: 
                st.markdown("<div style='margin-top: 30px;'></div>", unsafe_allow_html=True)
                ai_int = st.checkbox("✨ AI 보완", key="chk_int", value=False)
            with c4: 
                # 🚨 기본값 "1장" (index 1)
                sel_int = st.selectbox("슬라이드 배분", slide_options, index=1, key="sel_int")

        st.markdown("<br>", unsafe_allow_html=True)
        
        # --- 중단 (Row 2): AI 인사이트 & 뉴스 ---
        st.markdown("#### 🤖 AI 인사이트 & 뉴스")
        row2_col1, row2_col2 = st.columns(2)
        
        with row2_col1:
            ai_insight = st.text_area("3. AI Insight (핵심 요청사항이나 방향성을 입력)", height=150)
            c5, c6 = st.columns([1, 1])
            with c5: 
                st.markdown("<div style='margin-top: 30px;'></div>", unsafe_allow_html=True)
                ai_ins = st.checkbox("✨ AI 보완", key="chk_ins", value=True)
            with c6: 
                # 🚨 기본값 "3장" (index 3)
                sel_ins = st.selectbox("슬라이드 배분", slide_options, index=3, key="sel_ins")

        with row2_col2:
            ai_news = st.text_area("4. AI 뉴스", height=150)
            c7, c8 = st.columns([1, 1])
            with c7: 
                st.markdown("<div style='margin-top: 30px;'></div>", unsafe_allow_html=True)
                ai_news_chk = st.checkbox("✨ AI 보완", key="chk_news", value=True)
            with c8: 
                # 🚨 기본값 "1장" (index 1)
                sel_news = st.selectbox("슬라이드 배분", slide_options, index=1, key="sel_news")

        st.markdown("<br>", unsafe_allow_html=True)
        
        # --- 하단 (Row 3): 파일 업로더 ---
        st.markdown("##### 📂 AI Insight 참고 자료 첨부 (선택 사항)")
        uploaded_files = st.file_uploader("기획에 참고할 문서(PDF, Word, PPT, Excel)를 자유롭게 업로드하세요.", type=["pdf", "docx", "pptx", "xlsx"], accept_multiple_files=True)
        
        st.markdown("---")
        submit = st.form_submit_button("🚀 초안 및 디자인 추천 생성하기", use_container_width=True)
        
# 폼 제출 로직
        if submit:
            if not month: 
                st.warning("발행 월을 입력해주세요!")
            else:
                with st.spinner("첨부파일을 분석하고 초안과 맞춤형 PPT 디자인 테마를 기획 중입니다..."):
                    
                    # 🚨 [여기서부터 복구!] 첨부파일에서 텍스트를 뽑아내는 로직입니다.
                    extracted_text = ""
                    if uploaded_files:
                        for file in uploaded_files:
                            extracted_text += f"\n--- [{file.name}] 내용 ---\n"
                            extracted_text += extract_text_from_file(file)
                    # 🚨 [여기까지]
                    
                    # 이제 extracted_text가 정의되었으므로 에러가 나지 않습니다.
                    combined_ai_insight = f"사용자 방향성: {ai_insight}\n\n[참고 문서 내용]\n{extracted_text}" if extracted_text else ai_insight
                    
                    # 선택된 문자열을 숫자나 None으로 변환하는 함수
                    def parse_slide_cnt(val):
                        return None if val == "AI 자율" else int(val.replace("장", ""))
                    
                    # data 딕셔너리 완성
                    data = {
                        "month": month, 
                        "financial": financial, "ai_fin": ai_fin, "cnt_fin": parse_slide_cnt(sel_fin),
                        "internal": internal, "ai_int": ai_int, "cnt_int": parse_slide_cnt(sel_int),
                        "ai_insight": combined_ai_insight, "ai_ins": ai_ins, "cnt_ins": parse_slide_cnt(sel_ins),
                        "ai_news": ai_news, "ai_news_chk": ai_news_chk, "cnt_news": parse_slide_cnt(sel_news)
                    }
                    
                    # 초안 생성
                    st.session_state.draft_content = generate_draft(data)
                    st.session_state.teaser_content = generate_teaser(data) 
                    st.session_state.design_recommendation = get_design_recommendation(month, combined_ai_insight)
                    
                    st.session_state.month_title = month
                    st.session_state.step = 2
                    st.rerun()

# ---------------------------------------------------------
# Step 2: 초안 검토 및 수정 (직접 수정 기능 추가)
# ---------------------------------------------------------
elif st.session_state.step == 2:
    st.subheader("🧐 Step 2. Outline 초안 및 티저 직접 검토")
    st.info("아래 텍스트 박스에서 직접 내용을 수정하실 수 있습니다. 수정 후 우측 하단의 '저장 및 검토 완료'를 누르면 변경 사항이 PPT에 반영됩니다.")
    
    col_draft, col_teaser = st.columns([2, 1])
    with col_draft:
        st.markdown("### 📄 PPT 추출용 Outline 스크립트")
        edited_draft = st.text_area("초안 내용", value=st.session_state.draft_content, height=450, label_visibility="collapsed")
        
    with col_teaser:
            # 🚨 추가: 상단을 두 칸으로 나누어 제목과 버튼을 나란히 배치
            t_col1, t_col2 = st.columns([2, 1])
            with t_col1:
                st.markdown("### 🎨 티저 기획안")
            with t_col2:
                if st.button("🔄 새로 구상", use_container_width=True):
                    with st.spinner("수정된 스크립트로 다시 구상 중..."):
                        # 🚨 추가: 직접 수정한 edited_draft를 기반으로 다시 생성
                        st.session_state.teaser_content = regenerate_teaser_from_draft(edited_draft)
                        st.rerun()
                        
            edited_teaser = st.text_area("티저 내용", value=st.session_state.teaser_content, height=400, label_visibility="collapsed")
    
    st.markdown("---")
    st.markdown("#### 🤖 AI에게 추가 수정 요청하기 (선택 사항)")
    feedback = st.chat_input("채팅으로 수정할 내용을 입력하세요 (예: 3페이지 실적 부분은 조금 더 요약해줘)")
    
    if feedback:
        with st.chat_message("user"): st.write(feedback)
        with st.spinner("피드백 반영 중..."):
            st.session_state.draft_content = revise_draft(edited_draft, feedback)
            st.session_state.teaser_content = edited_teaser
            st.rerun()
            
    st.divider()
    col_btn1, col_btn2 = st.columns(2)
    with col_btn1:
        if st.button("처음부터 다시 입력하기", use_container_width=True):
            st.session_state.step = 1; st.rerun()
    with col_btn2:
        if st.button("✅ 저장 및 검토 완료. 출력 단계로 이동", type="primary", use_container_width=True):
            st.session_state.draft_content = edited_draft
            st.session_state.teaser_content = edited_teaser
            st.session_state.step = 3; st.rerun()

# ---------------------------------------------------------
# Step 3: 최종 출력
# ---------------------------------------------------------
elif st.session_state.step == 3:
    st.subheader("🎉 Step 3. 최종 디자인 및 PPT 생성")
    
    # 🚨 [수정] '페이지' 대신 '슬라이드' 번호를 찾고 +1을 해줍니다.
    if "final_full_text" not in st.session_state:
        draft_text = st.session_state.draft_content
        
        # '슬라이드 1', '슬라이드 8' 등에 있는 숫자만 모두 추출
        slide_numbers = re.findall(r'슬라이드\s*(\d+)', draft_text)
        
        if slide_numbers:
            last_slide = max(map(int, slide_numbers))
            next_slide = last_slide + 1
        else:
            next_slide = 99 # 번호를 못 찾을 경우 예비값
        
        # 툴들이 인식하기 쉽게 슬라이드 N. 포맷으로 티저 삽입
        teaser_title = f"\n\n==================================================\n슬라이드 {next_slide}. 티저 슬라이드 기획안 (표지/하이라이트용)\n==================================================\n"
        st.session_state.final_full_text = f"{draft_text}{teaser_title}{st.session_state.teaser_content}"
        
    report_title = f"expl'AI'n telink - {st.session_state.month_title}"
    
    st.success("✨ **AI가 제안하는 이번 달 맞춤형 디자인 테마 3선**")
    st.markdown(st.session_state.design_recommendation)
    
    st.divider()
    
    # --- 디자인 테마 원클릭 적용 및 직접 입력 버튼 ---
    if "selected_manus_style" not in st.session_state:
        st.session_state.selected_manus_style = ""

    prompts = re.findall(r'```(.*?)```', st.session_state.design_recommendation, re.DOTALL)
    
    if prompts:
        st.write("💡 **마음에 드는 테마를 클릭하거나, 직접 테마를 입력해 슬라이드 제작에 적용하세요!**")
        
        cols = st.columns(min(len(prompts) + 1, 4)) 
        
        for i, prompt_text in enumerate(prompts[:3]):
            with cols[i]:
                if st.button(f"🎨 추천 테마 {i+1} 적용", key=f"theme_btn_{i}", use_container_width=True):
                    st.session_state.selected_manus_style = prompt_text.strip()
                    st.rerun()
                    
        with cols[-1]:
            if st.button("✍️ 테마 직접 입력하기", type="primary", use_container_width=True):
                st.session_state.selected_manus_style = "" # 텍스트 박스 초기화
                st.rerun()
                
    st.divider()
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.write("📂 **1. 텍스트 내보내기 및 최종 편집**")
        st.info("아래 텍스트 박스에서 내용을 최종 수정한 후, 전체 복사(Ctrl+A, Ctrl+C)하여 붙여넣으세요. 수정한 내용은 PDF와 Manus 제작에 자동 반영됩니다.")
        
        edited_final_text = st.text_area(
            "최종 스크립트 수정", 
            value=st.session_state.final_full_text, 
            height=450,
            label_visibility="collapsed"
        )
        
        st.session_state.final_full_text = edited_final_text
        
        pdf_bytes, safe_name = create_professional_pdf(edited_final_text, report_title)
        st.download_button("📩 수정된 대본 PDF 다운로드", data=pdf_bytes, file_name=f"{safe_name}.pdf", use_container_width=True)
        
        if st.button("🔄 이전 단계로 돌아가기", use_container_width=True):
            # 이전 단계로 돌아갈 때 텍스트를 비워둬야 다시 넘어올 때 페이지 번호를 새로 계산함
            if "final_full_text" in st.session_state:
                del st.session_state["final_full_text"]
            st.session_state.step = 2; st.rerun()
            
    with col2:
        st.write("🎨 **시각화 자료**")
        
        if "manus_status" not in st.session_state:
            st.session_state.manus_status = "idle" 
        if "manus_url" not in st.session_state:
            st.session_state.manus_url = None

        # 🚨 [수정] st.expander 대신 st.container를 사용하여 중첩 에러를 방지합니다.
        with st.container(border=True):
            st.markdown("#### 📊 Manus 인포그래픽 슬라이드 제작")
            
            style_input = st.text_area(
                "디자인 테마 (버튼으로 선택하거나 직접 입력하세요)", 
                value=st.session_state.selected_manus_style,
                height=250,
                disabled=(st.session_state.manus_status == "processing")
            )

            if st.session_state.manus_status == "idle":
                if st.button("🚀 슬라이드 생성 시작", use_container_width=True):
                    st.session_state.selected_manus_style = style_input 
                    st.session_state.manus_status = "processing"
                    st.rerun()

            elif st.session_state.manus_status == "processing":
                st.button("⏳ 슬라이드 제작 중... 잠시만 기다려주세요", disabled=True, use_container_width=True)
                
                # 이제 컨테이너 안에 있으므로 st.status가 에러 없이 잘 작동합니다!
                with st.status("📊 Manus 에이전트 가동 중...", expanded=True) as status:
                    url, msg = create_manus_infographic(report_title, edited_final_text, style_input)
                    if url:
                        st.session_state.manus_url = url
                        st.session_state.manus_status = "completed"
                        status.update(label="✅ 제작 완료!", state="complete")
                        st.rerun()
                    else:
                        st.session_state.manus_status = "idle"
                        st.error(f"❌ 오류 발생: {msg}")
                        if st.button("🔄 다시 시도"): st.rerun()

            elif st.session_state.manus_status == "completed":
                st.link_button("📂 제작된 슬라이드 확인하기", st.session_state.manus_url, use_container_width=True, type="primary")
                if st.button("🆕 새로 만들기", use_container_width=True):
                    st.session_state.manus_status = "idle"
                    st.session_state.manus_url = None
                    st.rerun()